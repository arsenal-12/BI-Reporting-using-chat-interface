from flask import Flask, render_template, request, redirect, url_for,session,send_file
import requests
from flask_session import Session
from io import BytesIO,StringIO
from dotenv import load_dotenv
import os
import sqlite3
import google.generativeai as genai
import pandas as pd
import re
import json
from flask import jsonify
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import matplotlib.pyplot as plt
from matplotlib.backends.backend_agg import FigureCanvasAgg as FigureCanvas
from io import BytesIO
import base64
import tempfile
from pptx import Presentation
from pptx.util import Inches,Pt
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image
from reportlab.lib.styles import getSampleStyleSheet,ParagraphStyle
import matplotlib.pyplot as plt
import pandas as pd
import json
from io import BytesIO, StringIO
import tempfile
from flask import Flask, request, send_file
import markdown2
from pptx.enum.text import PP_ALIGN
import time
 
 
app = Flask(__name__)
 
# Load environment variables
load_dotenv()
# Set a secret key for session management
app.secret_key = os.getenv('FLASK_SECRET_KEY', 'your_secret_key_here')
# Load API key from environment variable
api_key = os.getenv('GOOGLE_API_KEY')
 
if api_key:
    genai.configure(api_key=api_key)
else:
    raise EnvironmentError("API key not found. Please set the GOOGLE_API_KEY environment variable.")
 
def get_gemini_response(question, prompt):
    combined_input = f"{prompt}\n\n{question}"
   
    try:
        response = genai.generate_text(
            model="models/text-bison-001",
            prompt=combined_input
        )
       
        generated_text = response.result
        sql_query = extract_sql_from_text(generated_text)
        chart_type = extract_chart_type_from_text(question)
        chart_config = extract_chart_config_from_text(generated_text)
       
        if sql_query:
            return sql_query, question, chart_type, None
        elif chart_config:
            return None, None, None, chart_config
        else:
            return None, None, None, None
    except Exception as e:
        print(f"Error in get_gemini_response: {e}")
        return None, None, None, None
 
 
def extract_sql_from_text(generated_text):
    patterns = [
        r'SELECT.*?FROM.*?;',
        r'INSERT.*?INTO.*?;',
        r'UPDATE.*?SET.*?;',
        r'DELETE.*?FROM.*?;',
        r'CREATE.*?;',
        r'DROP.*?;'
    ]
   
    for pattern in patterns:
        match = re.search(pattern, generated_text, re.DOTALL | re.IGNORECASE)
        if match:
            return match.group().strip()
   
    return None
 
def extract_chart_type_from_text(question):
    question = question.lower()
    if "pie" in question:
        return "pie"
    elif "bar" in question or "histogram" in question:
        return "bar"
    elif "line" in question:
        return "line"
    elif "scatter" in question:
        return "scatter"
    else:
        return None
 
def extract_chart_config_from_text(generated_text):
    try:
        chart_config_pattern = r'\{.*?\}'
        match = re.search(chart_config_pattern, generated_text, re.DOTALL)
        if match:
            chart_config_str = match.group()
            chart_config = json.loads(chart_config_str)
            return chart_config
        else:
            return None
    except json.JSONDecodeError as e:
        print(f"Error parsing chart config: {e}")
   
    return None
 
def read_sql_query(sql, db):
    try:
        conn = sqlite3.connect(db)
        df = pd.read_sql_query(sql, conn)
        conn.close()
        return df
    except sqlite3.Error as e:
        print(f"SQLite error: {e}")
        return None
 
def check_if_example_exists(prompt, question, sql_query):
    question_exists = question in prompt
    sql_exists = sql_query in prompt
    return question_exists and sql_exists
 
 
def generate_echarts(chart_data, chart_type="line", editable_colors=False):
    try:
        echarts_url = 'https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js'
        echarts_script = f'<script src="{echarts_url}"></script>'
       
        chart_json=json.dumps(chart_data)
        html_content = f'''
        <div>
            <label for="colorPicker">Select a color:</label>
            <input type="color" id="colorPicker">
        </div>
        <div>
            <form id="chartTitleForm">
                <label for="chartTitle">Chart Title:</label>
                <input type="text" id="chartTitle" name="chartTitle" value="Chart Title-1">
                <button type="button" onclick="updateChart()">Submit</button>
            </form>
        </div>
        <div id="chart" style="width: 800px; height: 600px;"></div>
        {echarts_script}
        <script type="text/javascript">
            var chartDom = document.getElementById('chart');
            var myChart = echarts.init(chartDom);
            var option = {json.dumps(chart_data)};
            myChart.setOption(option);
 
            document.getElementById('colorPicker').addEventListener('input', function() {{
                var color = this.value;
                option.series[0].itemStyle = {{ color: color }};
                myChart.setOption(option);
            }});
            document.getElementById('chartTitleForm').addEventListener('submit', function(event) {{
                event.preventDefault();
                updateChart();
            }});
            function updateChart() {{
                var title = document.getElementById('chartTitle').value;
                option.title = {{ text: title }};
                myChart.setOption(option);
                updateChartData();
            }}
            function updateChartData() {{
                var xhr = new XMLHttpRequest();
                xhr.open('POST', '/update_chart_data', true);
                xhr.setRequestHeader('Content-Type', 'application/json');
                xhr.send(JSON.stringify(option));
            }}
        </script>
        '''
        with open('chart_data.json','w') as f:
            json.dump(chart_data,f)
        return html_content
   
    except Exception as e:
        print(f"Error generating ECharts: {e}")
        return None
   
def get_next_example_number(prompt):
    examples = re.findall(r'Example \d+', prompt)
    if examples:
        example_numbers = [int(re.search(r'\d+', ex).group()) for ex in examples]
        next_example_number = max(example_numbers) + 1
    else:
        next_example_number = 1  # Start with 1 if no examples are found
    return next_example_number
 
def update_prompt_with_example(prompt, user_question, sql_query):
    next_example_number = get_next_example_number(prompt)
    new_example = f"Example {next_example_number}: {user_question}\nSQL command: `{sql_query}`\n\n"
   
    with open(prompt_file_path, "a") as file:
        file.write(new_example)
 
import uuid
 
def create_chat():
    chat_id = str(uuid.uuid4())
    return chat_id
 
def add_message_to_chat(chat_id, question, response):
    try:
        conn = sqlite3.connect("northwind.db")
        cursor = conn.cursor()
        cursor.execute("INSERT INTO chats (chat_id, question, response, last_activity) VALUES (?, ?, ?, datetime('now'))", (chat_id, question, response))
        conn.commit()
        conn.close()
    except sqlite3.Error as e:
        print(f"SQLite error:{e}")    
 
 
def get_chat_history(chat_id):
    try:
        conn = sqlite3.connect("northwind.db")
        cursor = conn.cursor()
        cursor.execute("SELECT question, response FROM chats WHERE chat_id = ?", (chat_id,))
        chat_history = cursor.fetchall()
        conn.close()
        return chat_history
    except sqlite3.Error as e:
        print(f"SQLite error: {e}")
        return []
 
def get_all_chat_sessions():
    try:
        conn = sqlite3.connect("northwind.db")
        cursor = conn.cursor()
        cursor.execute("SELECT chat_id, question, response FROM chats WHERE last_activity >= datetime('now', '-30 days')")
        sessions = cursor.fetchall()
        conn.close()
       
        return [{'chat_id': row[0], 'name': f"Chat {row[0][:8]}", 'question': row[1], 'response': row[2]} for row in sessions]
    except sqlite3.Error as e:
        print(f"SQLite error: {e}")
        return []  # Return an empty list or handle the error appropriately
    except Exception as e:
        print(f"Unexpected error: {e}")
        return []  # Handle other unexpected errors here
 
 
 
def generate_table_html(data):
    if data is None or data.empty:
        return "<p>No data available to display.</p>"
   
    table_html = '<table class="table table-striped">'
    table_html += '<thead><tr>'
   
    for column in data.columns:
        table_html += f'<th scope="col">{column}</th>'
    table_html += '</tr></thead>'
    table_html += '<tbody>'
   
    for _, row in data.iterrows():
        table_html += '<tr>'
        for value in row:
            table_html += f'<td>{value}</td>'
        table_html += '</tr>'
    table_html += '</tbody></table>'
   
    return table_html
 
prompt_file_path = "BI reporting/prompt.txt"
metadata_file = 'metadata.json'
# Function to update prompt.txt with user example
def update_prompt_with_example(prompt, user_question, sql_query):
    next_example_number = get_next_example_number(prompt)
    new_example = f"Example {next_example_number}: {user_question}\nSQL command: `{sql_query}`\n\n"
   
    with open("BI reporting/prompt.txt", "a") as file:
        file.write(new_example)
 
@app.route('/feedback/<chat_id>', methods=['POST'])
def feedback(chat_id):
    satisfied = request.form.get('satisfied')
    if satisfied in ['true', 'false']:
        action = request.form['satisfied']
       
        # Retrieve chat history and last question
        chat_history = get_chat_history(chat_id)
        last_question = chat_history[-1][0] if chat_history else None
 
        if action == 'true':
            sql_query, user_question, chart_type, chart_config = get_gemini_response(last_question, prompt)
            if sql_query and user_question:
                update_prompt_with_example(prompt, user_question, sql_query)
            return redirect(url_for('chat', chat_id=chat_id))
        else:
            return redirect(url_for('chat', chat_id=chat_id))
 
    return redirect(url_for('chat', chat_id=chat_id))
 
import re
 
 
 
# Function to determine the next example number
def get_next_example_number(prompt):
    examples = re.findall(r'Example \d+', prompt)
    if examples:
        example_numbers = [int(re.search(r'\d+', ex).group()) for ex in examples]
        next_example_number = max(example_numbers) + 1
    else:
        next_example_number = 1  # Start with 1 if no examples are found
    return next_example_number
 
 
# Example of reading initial prompt from file (adjust as per your needs)
with open("BI reporting/prompt.txt", "r") as file:
    prompt = file.read()
 
 
def get_metadata():
    if os.path.exists(metadata_file):
        with open(metadata_file, 'r') as file:
            return json.load(file)
    return None
 
def construct_prompt(metadata):
    if metadata:
        database_name = metadata.get('database', 'northwind.db')
        tables = metadata.get('tables', {})
       
        initial_prompt = f"You are an expert in converting English questions to SQL queries!\n"
        initial_prompt += f"The SQL database is named `{database_name}` and contains the following tables with their respective columns:\n\n"
 
        for table, info in tables.items():
            column_names = ', '.join(col['name'] for col in info['columns'])
            initial_prompt += f"{table}: {column_names}\n"
       
        initial_prompt += "\nFor example:\n\n"
        initial_prompt +="Example 1: How many customers are there?\n\n"
        initial_prompt +="SQL command: `SELECT COUNT(*) FROM customers;`\n\n"
 
        initial_prompt +="Example 2: List all orders placed on 2016-07-04.\n\n"
        initial_prompt +="SQL command: `SELECT * FROM orders WHERE OrderDate = '2016-07-04';\n\n"
 
        initial_prompt +="Example 3: Find the total sales amount for each category in 2016.\n\n"
        initial_prompt +="SQL command: `SELECT Categories.CategoryName, SUM([OrderDetails].UnitPrice * [OrderDetails].Quantity * (1 - [OrderDetails].Discount)) AS TotalSales FROM  [OrderDetails] INNER JOIN  Orders ON [OrderDetails].OrderID = Orders.OrderID INNER JOIN  Products ON [OrderDetails].ProductID = Products.ProductID INNER JOIN  Categories ON Products.CategoryID = Categories.CategoryID WHERE Orders.OrderDate BETWEEN '2016-07-04' AND '2016-07-09' GROUP BY Categories.CategoryName;`\n\n"
 
 
        initial_prompt +="Example 4: give me details of order with orderID 10248\n\n"
        initial_prompt +="SQL command: `SELECT * FROM orders WHERE OrderID = 10248;`\n\n"
 
        initial_prompt +="Example 5:Retrieve details of postal codes for region Western Europe.\n\n"
        initial_prompt +="SQL command: `SELECT PostalCode FROM Customers WHERE Region = 'Western Europe';`\n\n"
 
        initial_prompt +="Example 6:Retrieve products information for a specific CategoryID 1.\n\n"
        initial_prompt +="SQL command: `SELECT ProductName, UnitPrice, UnitsInStock FROM Products WHERE CategoryID = 1;`\n\n"
 
        initial_prompt +="Example 7:Retrieve employee details by their last name 'Davolio'.\n\n"
        initial_prompt +="SQL command: `SELECT FirstName, Title, HireDate FROM Employees WHERE LastName = 'Davolio';`\n\n"
 
        initial_prompt +="Example 8:Retrieve order details for a customer 'ALFKI'\n\n"
        initial_prompt +="SQL command: `SELECT OrderID, OrderDate, ShippedDate, ShipName, ShipCountry FROM Orders WHERE CustomerID = 'ALFKI';`\n\n"
 
        initial_prompt +="Example 9:Retrieve details of postal codes for a specific region Western Europe.\n\n"
        initial_prompt +="SQL command:`SELECT PostalCode FROM Customers WHERE Region = 'Western Europe';`\n\n"
 
        initial_prompt +="Example 10:Plot a line chart showing the total sales amount per month for the year 2016..\n\n"
        initial_prompt +="SQL command:`SELECT strftime('%m', Orders.OrderDate) AS Month,SUM([OrderDetails].UnitPrice * [OrderDetails].Quantity * (1 - [OrderDetails].Discount)) AS TotalSales FROM [OrderDetails] INNER JOIN Orders ON [OrderDetails].OrderID = Orders.OrderID INNER JOIN Products ON [OrderDetails].ProductID = Products.ProductID WHERE strftime('%Y', Orders.OrderDate) = '2016' GROUP BY Month;`\n\n"
 
        initial_prompt +="Example 11:Create a bar chart displaying the number of orders shipped by each shipper..\n\n"
        initial_prompt +="SQL command:`SELECT ShipperID, COUNT(*) AS NumberOfOrders FROM Shippers GROUP BY ShipperID;`\n\n"
 
        initial_prompt +="Example 12:Generate a pie chart showing the distribution of products by category..\n\n"
        initial_prompt +="SQL command:`SELECT CategoryName, COUNT(*) AS NumberOfProducts FROM Products JOIN Categories ON Products.CategoryID = Categories.CategoryID GROUP BY CategoryName;`\n\n"
 
        initial_prompt +="Example 13:Create a histogram showing the distribution of unit prices for products\n\n"
        initial_prompt +="SQL command:`SELECT UnitPrice FROM Products;`\n\n"
 
        initial_prompt +="Example 14:Plot a scatter plot showing the relationship between order quantity and unit price.\n\n"
        initial_prompt +="SQL command:`SELECT Quantity AS OrderQuantity, UnitPrice FROM [OrderDetails] GROUP BY Quantity, UnitPrice ORDER BY Quantity, UnitPrice;`\n\n"
 
        initial_prompt +="Example 15:Generate an area chart representing the total sales amount per year.\n\n"
        initial_prompt +="SQL command:`SELECT YEAR(OrderDate) AS Year,SUM(UnitPrice * Quantity * (1 - Discount)) AS TotalSales FROM Orders JOIN [OrderDetails] ON Orders.OrderID = [OrderDetails].OrderID GROUP BY YEAR(OrderDate) ORDER BY YEAR(OrderDate);`\n\n"
 
        initial_prompt +="Example 16:Create a box plot to visualize the distribution of order quantities.\n\n"
        initial_prompt +="SQL command:SELECT Quantity FROM [OrderDetails];`\n\n"
 
        initial_prompt +="Example 17: Plot a time series chart showing the evolution of total sales amount over time.\n\n"
        initial_prompt +="SQL command:`SELECT OrderDate, SUM(UnitPrice * Quantity * (1 - Discount)) AS TotalSales FROM Orders JOIN [OrderDetails] ON Orders.OrderID = [OrderDetails].OrderID GROUP BY OrderDate ORDER BY OrderDate;`\n\n"
 
 
        return initial_prompt
   
    return None
 
metadata = get_metadata()
 
if metadata is None or 'tables' not in metadata:
    raise ValueError("Metadata is not properly loaded or does not contain 'tables'.")
 
initial_prompt = construct_prompt(metadata)
if os.path.exists(prompt_file_path):
    with open(prompt_file_path, "w") as file:
        file.write(initial_prompt)
 
with open("BI reporting/prompt.txt", "r") as file:
    prompt = file.read()
 
@app.route('/', methods=['GET', 'POST'])
def index():
 
    return render_template('index.html', initial_prompt=prompt, metadata=metadata)
 
def update_prompt_with_example(prompt, user_question, sql_query):
    next_example_number = get_next_example_number(prompt)
    new_example = f"Example {next_example_number}: {user_question}\nSQL command: `{sql_query}`\n\n"
   
    with open(prompt_file_path, "a") as file:
        file.write(new_example)
# Flask routes
 
@app.route('/new_chat', methods=['GET', 'POST'])
def new_chat():
    if request.method == 'POST':
        session.clear()
        new_chat_id = create_chat()  # Create a new chat session
        session['chat_id'] = new_chat_id
        return redirect(url_for('chat', chat_id=new_chat_id))
 
    existing_sessions = get_all_chat_sessions()
    return render_template('new_chat.html', existing_sessions=existing_sessions)
 
@app.route('/update_chart_data', methods=['POST'])
def update_chart_data():
    chart_data = request.json
    with open('updated_chart_data.json', 'w') as f:
        json.dump(chart_data, f)
    return '', 204
 
 
 
def clean_title(query):
    # Remove leading words like "find", "show", etc., and the question mark
    cleaned_query = re.sub(r'^(find|show|retrieve|plot the graph|plot|what is)\s+', '', query, flags=re.IGNORECASE)
    cleaned_query = cleaned_query.rstrip('?')
    return cleaned_query.strip()
 
def create_pdf(title, table_html):
    # Load chart data from JSON
    with open('updated_chart_data.json', 'r') as f:
        chart_data = json.load(f)
 
    # Create matplotlib figure and plot chart based on chart_data
    fig, ax = plt.subplots(figsize=(10, 6))
    if chart_data['series'][0]['type'] == 'bar':
        ax.bar(chart_data['xAxis']['data'], chart_data['series'][0]['data'], color=chart_data['series'][0].get('itemStyle', {}).get('color'))
    elif chart_data['series'][0]['type'] == 'line':
        ax.plot(chart_data['xAxis']['data'], chart_data['series'][0]['data'], marker='o', color=chart_data['series'][0].get('itemStyle', {}).get('color'))
    elif chart_data['series'][0]['type'] == 'pie':
        ax.pie([d['value'] for d in chart_data['series'][0]['data']], labels=[d['name'] for d in chart_data['series'][0]['data']],
               colors=[d.get('itemStyle', {}).get('color') for d in chart_data['series'][0]['data']])
    ax.set_title(chart_data['title']['text'])
 
    # Save the chart as a PNG image
    img_buffer = BytesIO()
    fig.savefig(img_buffer, format='png')
    img_buffer.seek(0)
 
    # Create a PDF using ReportLab
    pdf_output = BytesIO()
    doc = SimpleDocTemplate(pdf_output, pagesize=letter)
    elements = []
 
    # Add Title to PDF
    styles = getSampleStyleSheet()
    title_paragraph = Paragraph(title, styles['Title'])
    elements.append(title_paragraph)
    elements.append(Spacer(1, 12))
 
    # Add Chart Image to PDF
    img_buffer.seek(0)
    img = Image(img_buffer, width=400, height=250)
    elements.append(img)
    elements.append(Spacer(1, 12))
 
    # Add Table Data to PDF
    if isinstance(table_html, str):
        df = pd.read_html(StringIO(table_html))[0]
        data = [df.columns.values.tolist()] + df.values.tolist()
 
        # Create Table
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
 
        elements.append(table)
        elements.append(Spacer(1, 12))
 
        # Generate summary using AI
        summary_text = generate_summary(chart_data)
       
        # Add breaks between sections
        summary_text = summary_text.replace("Chart Title:", "Chart Title:\n").replace("Data:", "\nData:\n").replace("Summary:", "\nSummary:\n")
       
        # Define custom paragraph style for summary
        summary_style = ParagraphStyle(
            'SummaryStyle',
            parent=styles['BodyText'],
            fontSize=10,
            spaceAfter=12,
            leftIndent=10,
            rightIndent=10
        )
       
        # Convert markdown summary to HTML
        html_summary = markdown2.markdown(summary_text)
       
        # Create a paragraph with the formatted summary
        summary_paragraph = Paragraph(html_summary, summary_style)
        elements.append(summary_paragraph)
 
    # Build PDF
    doc.build(elements)
 
    pdf_output.seek(0)
    plt.close(fig)  # Close the Matplotlib figure to release resources
 
    return pdf_output
 
 
 
@app.route('/download_pdf', methods=['GET'])
def download_pdf():
    try:
        query = request.args.get('title')
        title = clean_title(query)
        table_html = request.args.get('table_html')
 
        pdf_output = create_pdf(title, table_html)
 
        # Save PDF content to a temporary file
        temp_pdf = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        temp_pdf.write(pdf_output.getvalue())
        temp_pdf.close()
 
        # Return the temporary file as an attachment
        return send_file(temp_pdf.name, as_attachment=True, download_name='report.pdf')
 
    except Exception as e:
        print(f"Error in creating or downloading PDF: {str(e)}")
        return str(e), 500
 
def clean_title(query):
    # Remove leading words like "find", "show", etc., and the question mark
    cleaned_query = re.sub(r'^(find|show|what is |plot the graph for| retrive|plot )\s+', '', query, flags=re.IGNORECASE)
    cleaned_query = cleaned_query.rstrip('?')
    return cleaned_query.strip()
 
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))
 
def get_gemini_response_chart(chart_json, prompt):
    model = genai.GenerativeModel('gemini-pro')
    response = model.generate_content([prompt, chart_json])
    return response.text
 
def generate_summary(chart_data):
    # Convert chart data to a summary using Gemini AI
    chart_json = json.dumps(chart_data)
    prompt = f"""
You are an expert data analyst!
 
You are given the following chart data in JSON format: {chart_json}
 
Please provide a detailed summary of the data presented in this chart.
"""
    summary_text = get_gemini_response_chart(chart_json, prompt)  # Get the summary text from Gemini AI
 
    # Format the summary according to specified structure
    summary = f"""
Chart Summary:
 
The chart presents data on the sales of different product categories within a specific industry. Each bar represents the total sales value for a particular category.
 
Key Findings:
 
{summary_text.strip()}  # Assuming summary_text contains bullet points as per your example
 
Overall Trend:
 
The chart indicates a downward trend in sales values across most categories from "Beverages" to "Seafood." This trend suggests that while "Beverages" and "Condiments" are driving sales, other categories may be experiencing reduced demand or market share loss
"""
    return summary.strip()  # Strip any leading/trailing whitespace before returning
 
def create_ppt(title, table_html):
    # Load chart data from JSON
    with open('updated_chart_data.json', 'r') as f:
        chart_data = json.load(f)
 
    # Create matplotlib figure and plot chart based on chart_data
    fig, ax = plt.subplots(figsize=(10, 6))
    if chart_data['series'][0]['type'] == 'bar':
        ax.bar(chart_data['xAxis']['data'], chart_data['series'][0]['data'], color=chart_data['series'][0].get('itemStyle', {}).get('color'))
    elif chart_data['series'][0]['type'] == 'line':
        ax.plot(chart_data['xAxis']['data'], chart_data['series'][0]['data'], marker='o', color=chart_data['series'][0].get('itemStyle', {}).get('color'))
    elif chart_data['series'][0]['type'] == 'pie':
        ax.pie([d['value'] for d in chart_data['series'][0]['data']], labels=[d['name'] for d in chart_data['series'][0]['data']],
               colors=[d.get('itemStyle', {}).get('color') for d in chart_data['series'][0]['data']])
    ax.set_title(chart_data['title']['text'])
 
    # Save the chart as a PNG image
    img_buffer = BytesIO()
    fig.savefig(img_buffer, format='png')
    img_buffer.seek(0)
 
    # Create a PowerPoint presentation
    prs = Presentation()
 
    # Slide 1 - Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    title_shape.text = title
 
    # Slide 2 - Chart slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1)
    pic = slide.shapes.add_picture(img_buffer, left, top, height=Inches(5))
 
    # Slide 3 - Table slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(1)
    top = Inches(1)
    if isinstance(table_html, str):
        df = pd.read_html(StringIO(table_html))[0]
        rows, cols = df.shape
        table = slide.shapes.add_table(rows + 1, cols, left, top, Inches(8), Inches(4)).table
       
        # Adding column names
        for col_idx, col_name in enumerate(df.columns):
            table.cell(0, col_idx).text = col_name
 
        # Adding the data
        for row_idx in range(rows):
            for col_idx in range(cols):
                table.cell(row_idx + 1, col_idx).text = str(df.iat[row_idx, col_idx])
 
    # Slide 4 - Summary slide
    summary_text = generate_summary(chart_data)
    
 
        # Function to add text to slide and handle overflow
    def add_text_to_slide(slide, text, left, top, width, height):
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
 
        paragraphs = text.split("\n")
        for i, para in enumerate(paragraphs):
            p = text_frame.add_paragraph()
            p.text = para.strip()
            if para.startswith("***") and para.endswith("**"):
                p.text = "*"+ para[3:-2]
                p.font.bold = True
            elif para.startswith("**") and para.endswith("**"):
                p.text= para[2:-2]
                p.font.bold = True
            else:
                p.text=para.strip()

            p.font.size = Pt(14)
            if i == 0:
                p.font.bold = True
                p.font.size = Pt(16)
           
            # Check if the text frame's height exceeds the textbox height
            if text_frame.text.count('\n') * Pt(14).pt > height.pt:
                # If it does, remove the paragraph from the text frame and return the remaining text
                text_frame._element.remove(p._element)
                remaining_text = "\n".join(paragraphs[i:])
                return remaining_text
       
        return None
 
    # Add summary text to slide(s)
    summary_text = generate_summary(chart_data)
    slide_layout = prs.slide_layouts[5]  # Blank layout
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(4)
 
    remaining_text = summary_text
    while remaining_text:
        slide = prs.slides.add_slide(slide_layout)
        remaining_text = add_text_to_slide(slide, remaining_text, left, top, width, height)
 
    # Slide - Thank you slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(3)
    top = Inches(3)
    textbox = slide.shapes.add_textbox(left, top, Inches(4), Inches(2))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = "Thank You"
    p.font.bold = True
    p.font.size = Pt(32)
    p.alignment = PP_ALIGN.CENTER
 
 
    # Save PPT to BytesIO object
    ppt_output = BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
 
    plt.close(fig)  # Close the Matplotlib figure to release resources
 
    return ppt_output
 
def create_pdf(title, table_html):
    # Load chart data from JSON
    with open('updated_chart_data.json', 'r') as f:
        chart_data = json.load(f)
 
    # Create matplotlib figure and plot chart based on chart_data
    fig, ax = plt.subplots(figsize=(10, 6))
    if chart_data['series'][0]['type'] == 'bar':
        ax.bar(chart_data['xAxis']['data'], chart_data['series'][0]['data'],
               color=chart_data['series'][0].get('itemStyle', {}).get('color'))
    elif chart_data['series'][0]['type'] == 'line':
        ax.plot(chart_data['xAxis']['data'], chart_data['series'][0]['data'], marker='o',
                color=chart_data['series'][0].get('itemStyle', {}).get('color'))
    elif chart_data['series'][0]['type'] == 'pie':
        ax.pie([d['value'] for d in chart_data['series'][0]['data']],
               labels=[d['name'] for d in chart_data['series'][0]['data']],
               colors=[d.get('itemStyle', {}).get('color') for d in chart_data['series'][0]['data']])
    ax.set_title(chart_data['title']['text'])
 
    # Save the chart as a PNG image
    img_buffer = BytesIO()
    fig.savefig(img_buffer, format='png')
    img_buffer.seek(0)
 
    # Create a PDF using ReportLab
    pdf_output = BytesIO()
    doc = SimpleDocTemplate(pdf_output, pagesize=letter)
    elements = []
 
    # Add Title to PDF
    styles = getSampleStyleSheet()
    title_paragraph = Paragraph(title, styles['Title'])
    elements.append(title_paragraph)
    elements.append(Spacer(1, 12))
 
    # Add Chart Image to PDF
    img_buffer.seek(0)
    img = Image(img_buffer, width=400, height=250)
    elements.append(img)
    elements.append(Spacer(1, 12))
 
    # Add Table Data to PDF
    if isinstance(table_html, str):
        df = pd.read_html(StringIO(table_html))[0]
        data = [df.columns.values.tolist()] + df.values.tolist()
 
        # Create Table
        table = Table(data)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
            ('GRID', (0, 0), (-1, -1), 1, colors.black),
        ]))
 
        elements.append(table)
        elements.append(Spacer(1, 12))
 
    # Generate summary using AI
    summary_text = generate_summary(chart_data)
 
    # Add breaks between sections
    summary_text = summary_text.replace("Chart Title:", "\nChart Title:\n").replace("Data:", "\nData:\n").replace("Summary:", "\nSummary:\n")
 
        # Split summary into bullet points
    summary_points = [line.strip() for line in summary_text.splitlines() if line.strip()]

    bold_style=ParagraphStyle('Bold',parent=styles['BodyText'],fontName = 'Helvetica-Bold')
 
# Add summary points to PDF as bullets
    for point in summary_points:
     if point.startswith("***") and point.endswith("**"):
            point_text = "*" + point[3:-2]
            elements.append(Paragraph(point_text, bold_style))
     elif point.startswith("**") and point.endswith("**"):
            point_text = point[2:-2]
            elements.append(Paragraph(point_text, bold_style))
     else:
            elements.append(Paragraph(point, styles['BodyText']))
     elements.append(Spacer(1, 6))
 
   
    # Build PDF
    doc.build(elements)
 
    pdf_output.seek(0)
    plt.close(fig)  # Close the Matplotlib figure to release resources
 
    return pdf_output
 
@app.route('/download_ppt', methods=['GET'])
def download_ppt():
    try:
        query = request.args.get('title')
        title = clean_title(query)
        table_html = request.args.get('table_html')
 
        ppt_output = create_ppt(title, table_html)
 
        # Save PPT content to a temporary file
        temp_ppt = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        temp_ppt.write(ppt_output.getvalue())
        temp_ppt.close()
 
        # Return the temporary file as an attachment
        return send_file(temp_ppt.name, as_attachment=True, download_name='report.pptx')
 
    except Exception as e:
        print(f"Error in creating or downloading PPT: {str(e)}")
        return str(e), 500
 
 
@app.route('/chat/<chat_id>', methods=['GET', 'POST'])
def chat(chat_id):
    if 'chat_id' not in session and session['chat_id'] != chat_id:
 
        return redirect(url_for('new_chat'))
   
    if request.method == 'POST':
        try:
            question = request.form['question']
            sql_query, user_question, chart_type, chart_config = get_gemini_response(question, prompt)
 
            if sql_query and user_question:
                example_exists = check_if_example_exists(prompt, user_question, sql_query)
                data = read_sql_query(sql_query, "northwind.db")
                chart_html = ""
                table_html = ""
                graph_html=""
                editable_colors = False
 
                if data is not None:
                    try:
                        if data.shape == (1, 1):
                            single_value = data.iloc[0, 0]
                            table_html = f"<p>Result: {single_value}</p>"
                        else:
                            if not chart_type:
                                if len(data.columns) == 2:
                                    chart_type = "bar"
                                elif len(data.columns) > 2:
                                    chart_type = "line"
                                elif len(data.columns) == 1:
                                    chart_type = "pie"
                                else:
                                    chart_type = "pie"
 
                            if chart_type == "pie":
                                chart_data = {
                                    'title': {'text': 'Chart Title'},
                                    'series': [{
                                        'type': 'pie',
                                        'data': [{'value': int(v), 'name': k} for k, v in zip(data.iloc[:, 0], data.iloc[:, 1])]
                                    }]
                                }
                            elif chart_type == "scatter":
                                chart_data = {
                                    'title': {'text': 'Chart Title'},
                                    'xAxis': {
                                        'type': 'value'
                                    },
                                    'yAxis': {
                                        'type': 'value'
                                    },
                                    'series': [{
                                        'data': [[x, y] for x, y in zip(data.iloc[:, 0], data.iloc[:, 1])],
                                        'type': 'scatter'
                                    }]
                                }
                            else:
                                chart_data = {
                                    'title': {'text': 'Chart Title'},
                                    'xAxis': {
                                        'type': 'category',
                                        'data': data.iloc[:, 0].tolist() if not data.empty else []
                                    },
                                    'yAxis': {
                                        'type': 'value'
                                    },
                                    'series': [{
                                        'data': data.iloc[:, 1].tolist() if not data.empty else [],
                                        'type': chart_type
                                    }]
                                }
 
                            editable_colors = chart_type in ["bar", "line", "pie"]
                            chart_html = generate_echarts(chart_data, chart_type, editable_colors=True)
                            table_html = generate_table_html(data)
                    except IndexError as e:
                        print(f"IndexError: {e}")
                        chart_html = "<p>Data format issue: The data does not contain the expected number of columns or rows.</p>"
                        table_html = "<p>No table data to display due to data format issue.</p>"
                    except Exception as e:
                        print(f"Unexpected Error: {e}")
                        chart_html = "<p>Unexpected error occurred while generating the chart.</p>"
                        table_html = "<p>Unexpected error occurred while generating the table.</p>"
 
                    add_message_to_chat(chat_id, question, sql_query)
                    chat_history = get_chat_history(chat_id)
                    show_feedback = not example_exists
                    return render_template('chat.html',question=question, chat_id=chat_id, chat_history=chat_history, chart_html=chart_html, table_html=table_html, show_feedback=show_feedback)
                else:
                    chat_history = get_chat_history(chat_id)
                    return render_template('chat.html', chat_id=chat_id, chat_history=chat_history, message="No data returned or an error occurred.")
        except IndexError as e:
            chat_history = get_chat_history(chat_id)
            return render_template('chat.html', chat_id=chat_id, chat_history=chat_history, message=f"Index error: {str(e)}")
        except Exception as e:
            chat_history = get_chat_history(chat_id)
            return render_template('chat.html', chat_id=chat_id, chat_history=chat_history, message=f"An error occurred: {str(e)}")
 
    chat_history = get_chat_history(chat_id)
    return render_template('chat.html', chat_id=chat_id, chat_history=chat_history)
 
 
# Function to get Gemini response based on chart data and query
def get_gemini_response_chart1(chart_json, chart_query, db_path):
    prompt = f"""
    You are an expert data analyst!\n
    You are given the following chart data in JSON format: {chart_json}\n
    The user has asked the following question about the chart: {chart_query}\n
    The data is sourced from the Northwind database located at: {db_path}\n
    Please provide a detailed response to the user's question.
    If the user asks you to plot the graph for their queries then return output in JSON format.
    """
    # Call function to get Gemini response
    response = genai.generate_text(
        model="models/text-bison-001",
        prompt=prompt,
        temperature=0.7,
        candidate_count=1,
        top_k=40,
        top_p=0.95,
        max_output_tokens=1024,
        stop_sequences=[]
    )
   
    return response.result
 
 
def extract_json_from_response(response):
    # Correct the regex pattern to use triple backticks
    json_code = re.search(r'```json\s*(\{.*?\})\s*```', response, re.DOTALL)
    if json_code:
        print(f"Matched JSON code:\n{json_code.group(1)}")  # Debug print to see the matched JSON code
        # Replace single quotes with double quotes
        json_str = json_code.group(1).replace("'", '"')
        try:
            parsed_json = json.loads(json_str)
            return parsed_json
        except json.JSONDecodeError as e:
            print(f"JSON decoding error: {e}")
            return None
 

def generate_updatedgraph(chart_data, chart_type="bar", editable_colors=False):
    try:
        echarts_url = 'https://cdn.jsdelivr.net/npm/echarts@5/dist/echarts.min.js'
        echarts_script = f'<script src="{echarts_url}"></script>'
        chart_id = f'chart_{int(time.time() * 1000)}'  # Unique ID based on timestamp

        html_content = f'''
        <div id="{chart_id}_container">
            <div>
                <label for="{chart_id}_colorPicker">Select a color:</label>
                <input type="color" id="{chart_id}_colorPicker">
            </div>
            <div>
                <form id="{chart_id}_chartTitleForm">
                    <label for="{chart_id}_chartTitle">Chart Title:</label>
                    <input type="text" id="{chart_id}_chartTitle" name="chartTitle" value="{chart_data.get('title', {}).get('text', 'Chart Title')}">
                    <button type="button" onclick="updateChart('{chart_id}')">Submit</button>
                </form>
            </div>
            <div id="{chart_id}" style="width: 800px; height: 600px;"></div>
            {echarts_script}
            <script type="text/javascript">
                var chartDom = document.getElementById('{chart_id}');
                var myChart = echarts.init(chartDom);
                var option = {json.dumps(chart_data)};
                myChart.setOption(option);

                document.getElementById('{chart_id}_colorPicker').addEventListener('input', function() {{
                    var color = this.value;
                    option.series[0].itemStyle = {{ color: color }};
                    myChart.setOption(option);
                }});
                document.getElementById('{chart_id}_chartTitleForm').addEventListener('submit', function(event) {{
                    event.preventDefault();
                    updateChart('{chart_id}');
                }});
                function updateChart(chartId) {{
                    var title = document.getElementById(chartId + '_chartTitle').value;
                    option.title = {{ text: title }}; 
                    myChart.setOption(option);
                    updateChartData();
                }}
                function updateChartData() {{
                    var xhr = new XMLHttpRequest();
                    xhr.open('POST', '/update_chart_data', true);
                    xhr.setRequestHeader('Content-Type', 'application/json');
                    xhr.send(JSON.stringify(option));
                }}
            </script>
        </div>
        '''
       
        return html_content
   
    except Exception as e:
        print(f"Error generating ECharts: {e}")
        return None


# Route to handle chart response generation
@app.route('/generate_chart_response', methods=['POST'])
def generate_chart_response():
    try:
        data = request.json
        chart_query = data.get('chart_query')  # Extract chart query from the request
 
        # Read the chart data from the JSON file
        with open('chart_data.json', 'r') as file:
            chart_json = json.load(file)
 
        db_path = 'northwind.db'  # Replace with your actual db_path
 
        # Call function to get Gemini response
        gemini_response = get_gemini_response_chart1(chart_json, chart_query, db_path)
 
        # Extract JSON from Gemini response
        extracted_json = extract_json_from_response(gemini_response)
 
        if extracted_json:
            # Save the extracted JSON as graph.json
            with open('graph.json', 'w') as file:
                json.dump(extracted_json, file)
 
            # Read the saved JSON data from graph.json
            with open('graph.json', 'r') as file:
                graph_data = json.load(file)
 
            # Generate the ECharts graph
            graph_html = generate_updatedgraph(graph_data)
            print(graph_html)
 
            if graph_html:
                return jsonify({'success': True, 'response': gemini_response, 'graph': graph_html})
            else:
                return jsonify({'success': False, 'error': 'Failed to generate chart HTML'})
 
        else:
            return jsonify({'success': False, 'error': 'Could not extract JSON from Gemini response'})
 
    except Exception as e:
        return jsonify({'success': False, 'error': f'Exception occurred: {str(e)}'})
 
if __name__ == '__main__':
    app.run(debug=True)